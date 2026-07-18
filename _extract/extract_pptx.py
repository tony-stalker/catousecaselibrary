#!/usr/bin/env python3
"""Extract text, notes, and images from all PPTX decks into a digest."""
import json, os, re, sys, zipfile
from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

ROOT = Path("/Users/tonystalker/Documents/claude/usecaselibrary")
OUT = ROOT / "_extract"
OUT.mkdir(exist_ok=True)
MEDIA = OUT / "media"
MEDIA.mkdir(exist_ok=True)

def slug(s):
    s = re.sub(r"[^A-Za-z0-9]+", "-", s).strip("-").lower()
    return s[:60]

def shape_text(shape, depth=0):
    items = []
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for sub in shape.shapes:
            items.extend(shape_text(sub, depth + 1))
        return items
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = "".join(run.text for run in para.runs).strip()
            if t:
                items.append({"lvl": para.level, "text": t})
    if getattr(shape, "has_table", False) and shape.has_table:
        rows = []
        for r in shape.table.rows:
            rows.append([c.text.strip() for c in r.cells])
        items.append({"table": rows})
    return items

def extract_deck(pptx_path: Path):
    deck = {"file": str(pptx_path.relative_to(ROOT)), "slides": []}
    try:
        prs = Presentation(str(pptx_path))
    except Exception as e:
        deck["error"] = str(e)
        return deck
    dslug = slug(pptx_path.stem)
    for i, slide in enumerate(prs.slides, 1):
        s = {"n": i, "title": "", "content": [], "notes": "", "images": []}
        try:
            if slide.shapes.title is not None and slide.shapes.title.text.strip():
                s["title"] = slide.shapes.title.text.strip()
        except Exception:
            pass
        for shape in slide.shapes:
            try:
                if shape == slide.shapes.title:
                    continue
            except Exception:
                pass
            s["content"].extend(shape_text(shape))
            # images
            def grab_images(sh):
                if sh.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for sub in sh.shapes:
                        grab_images(sub)
                elif sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    try:
                        img = sh.image
                        ext = img.ext
                        name = f"{dslug}_s{i:02d}_{len(s['images'])}.{ext}"
                        p = MEDIA / name
                        if not p.exists():
                            p.write_bytes(img.blob)
                        s["images"].append({"file": name, "size": len(img.blob),
                                            "w_emu": sh.width, "h_emu": sh.height})
                    except Exception as e:
                        s["images"].append({"error": str(e)})
            grab_images(shape)
        try:
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
                s["notes"] = slide.notes_slide.notes_text_frame.text.strip()
        except Exception:
            pass
        deck["slides"].append(s)
    return deck

decks = []
for pptx in sorted(ROOT.rglob("*.pptx")):
    if "_extract" in str(pptx):
        continue
    print(f"Extracting: {pptx.name}", file=sys.stderr)
    decks.append(extract_deck(pptx))

(OUT / "digest.json").write_text(json.dumps(decks, indent=1))

# Also write a readable markdown digest per deck
for deck in decks:
    dslug = slug(Path(deck["file"]).stem)
    lines = [f"# {deck['file']}", ""]
    if "error" in deck:
        lines.append(f"ERROR: {deck['error']}")
    for s in deck.get("slides", []):
        lines.append(f"\n## Slide {s['n']}: {s['title'] or '(no title)'}")
        for item in s["content"]:
            if "table" in item:
                for row in item["table"]:
                    lines.append("| " + " | ".join(row) + " |")
            else:
                lines.append("  " * item["lvl"] + "- " + item["text"])
        if s["images"]:
            imgs = ", ".join(im.get("file", "ERR") for im in s["images"])
            lines.append(f"[images: {imgs}]")
        if s["notes"]:
            lines.append(f"NOTES: {s['notes']}")
    (OUT / f"{dslug}.md").write_text("\n".join(lines))

print(f"Done: {len(decks)} decks")
for d in decks:
    print(f"  {d['file']}: {len(d.get('slides', []))} slides" + (" ERROR" if "error" in d else ""))
